import sys
import re
from pptx import Presentation
from pathlib import Path

def insert_space_between_japanese_and_ascii(text):
    # 全角かっこを半角かっこに変換
    text = text.replace('（', '(').replace('）', ')')
    # 日本語（ひらがな・カタカナ・漢字）とASCII英数字の間にスペースを挿入
    text = re.sub(r'([A-Za-z0-9])([\u3040-\u30ff\u3400-\u4dbf\u4e00-\u9fff])', r'\1 \2', text)
    text = re.sub(r'([\u3040-\u30ff\u3400-\u4dbf\u4e00-\u9fff])([A-Za-z0-9])', r'\1 \2', text)
    return text

def process_pptx_file(input_path, output_path=None):
    prs = Presentation(input_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                # 1. paragraph 全体文字列を取得
                original_runs = paragraph.runs
                full_text = ''.join(r.text for r in original_runs)

                # 2. 半角スペースを挿入
                new_text = insert_space_between_japanese_and_ascii(full_text)

                # 3. 各 run の長さで分割して上書き
                lengths = [len(r.text) for r in original_runs]
                split_texts, idx = [], 0
                for ln in lengths:
                    split_texts.append(new_text[idx:idx+ln])
                    idx += ln
                if idx < len(new_text) and split_texts:
                    split_texts[-1] += new_text[idx:]     # 余りは最後に追加

                # 4. 上書き（run 数は変えないので書式は維持）
                for run, txt in zip(original_runs, split_texts):
                    run.text = txt

    if output_path is None:
        output_path = Path(input_path).with_name(
            Path(input_path).stem + "_spaced.pptx"
        )
    prs.save(output_path)
    print(f"✅ 保存しました: {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("使い方: python spaceNinja.py <input.pptx>")
        sys.exit(1)

    input_file = sys.argv[1]
    process_pptx_file(input_file)
